"""DOCX, PPTX, and XLSX accessibility checks."""
from pathlib import Path

from accessiflow.models import AccessibilityIssue, Severity
from accessiflow.checks.base import AccessibilityCheck
from accessiflow.checks.registry import register_check


@register_check
class DocxImagesMissingAlt(AccessibilityCheck):
    check_id = "docx-images-missing-alt"
    title = "DOCX images without alt text"
    description = "Images in Word documents should have alt text"
    wcag_criterion = "1.1.1"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".docx",):
            return []
        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            import xml.etree.ElementTree as ET

            doc = Document(str(file_path))
            issues = []

            # Check inline images for alt text
            nsmap = {
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }

            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    drawing_elements = run._element.findall('.//wp:inline', nsmap)
                    drawing_elements += run._element.findall('.//wp:anchor', nsmap)
                    for drawing in drawing_elements:
                        docPr = drawing.find('.//wp:docPr', nsmap)
                        if docPr is None:
                            # Try without namespace
                            for child in drawing.iter():
                                if child.tag.endswith('}docPr') or child.tag == 'docPr':
                                    docPr = child
                                    break

                        alt_text = ""
                        if docPr is not None:
                            alt_text = docPr.get('descr', '') or docPr.get('title', '')

                        if not alt_text.strip():
                            name = docPr.get('name', 'unknown') if docPr is not None else 'unknown'
                            issues.append(self._make_issue(
                                f"Image '{name}' in DOCX has no alt text",
                                severity=Severity.SERIOUS,
                            ))

            return issues
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze DOCX {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]


@register_check
class PptxSlidesMissingTitles(AccessibilityCheck):
    check_id = "pptx-slides-missing-titles"
    title = "PPTX slides without titles"
    description = "PowerPoint slides should have titles for navigation"
    wcag_criterion = "2.4.2"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".pptx",):
            return []
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            issues = []

            for i, slide in enumerate(prs.slides, 1):
                has_title = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.shape_id == slide.shapes.title_shape_id if hasattr(slide.shapes, 'title_shape_id') else False:
                        has_title = True
                        break
                    if hasattr(shape, "name") and "title" in shape.name.lower():
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_title = True
                            break

                # Simpler check: just look for slide.shapes.title
                if not has_title and slide.shapes.title is not None:
                    if slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text.strip():
                        has_title = True

                if not has_title:
                    issues.append(self._make_issue(
                        f"Slide {i} has no title",
                        severity=Severity.SERIOUS,
                    ))

            return issues
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze PPTX {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]


@register_check
class PptxImagesMissingAlt(AccessibilityCheck):
    check_id = "pptx-images-missing-alt"
    title = "PPTX images without alt text"
    description = "Images in PowerPoint presentations should have alt text"
    wcag_criterion = "1.1.1"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".pptx",):
            return []
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            issues = []

            for i, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        alt_text = ""
                        if hasattr(shape, "image"):
                            # Check alt text in shape properties
                            desc = shape._element.find(
                                './/{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr'
                            )
                            if desc is None:
                                for child in shape._element.iter():
                                    if child.tag.endswith('}cNvPr'):
                                        alt_text = child.get('descr', '') or child.get('title', '')
                                        break
                            else:
                                alt_text = desc.get('descr', '') or desc.get('title', '')

                        if not alt_text.strip():
                            name = shape.name if hasattr(shape, 'name') else 'unknown'
                            issues.append(self._make_issue(
                                f"Slide {i}: Image '{name}' has no alt text",
                                severity=Severity.SERIOUS,
                            ))

            return issues
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze PPTX images {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]


@register_check
class XlsxAccessibilityCheck(AccessibilityCheck):
    check_id = "xlsx-missing-sheet-titles"
    title = "XLSX sheets without descriptive names"
    description = "Spreadsheet sheets should have descriptive names, not defaults like Sheet1"
    wcag_criterion = "2.4.2"

    # Default sheet name patterns across locales
    _DEFAULT_PATTERNS = {"sheet", "hoja", "feuille", "blatt", "foglio", "planilha"}

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".xlsx", ".xls"):
            return []
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            issues = []

            for sheet_name in wb.sheetnames:
                # Check for generic/default names like "Sheet1", "Sheet2"
                name_lower = sheet_name.lower().strip()
                is_default = any(
                    name_lower.startswith(p) and (
                        name_lower[len(p):].isdigit() or name_lower == p
                    )
                    for p in self._DEFAULT_PATTERNS
                )
                if is_default:
                    issues.append(self._make_issue(
                        f"Sheet '{sheet_name}' has a default/non-descriptive name",
                        severity=Severity.MINOR,
                    ))

            # Check for merged cells (accessibility concern for screen readers)
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                if hasattr(ws, 'merged_cells') and ws.merged_cells:
                    merge_count = len(list(ws.merged_cells.ranges))
                    if merge_count > 0:
                        issues.append(self._make_issue(
                            f"Sheet '{sheet_name}' has {merge_count} merged cell region(s) "
                            f"which can cause issues for screen readers",
                            severity=Severity.MODERATE,
                        ))

            wb.close()
            return issues
        except ImportError:
            return [self._make_issue(
                f"openpyxl not installed — cannot check XLSX files",
                severity=Severity.MINOR,
            )]
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze XLSX {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]


@register_check
class XlsxImagesMissingAlt(AccessibilityCheck):
    check_id = "xlsx-images-missing-alt"
    title = "XLSX images without alt text"
    description = "Images in spreadsheets should have alt text"
    wcag_criterion = "1.1.1"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".xlsx", ".xls"):
            return []
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=False, data_only=True)
            issues = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for image in ws._images:
                    # openpyxl Image objects don't directly expose alt text easily
                    # but we can flag images for review
                    desc = ""
                    if hasattr(image, 'description'):
                        desc = image.description or ""
                    if not desc.strip():
                        issues.append(self._make_issue(
                            f"Sheet '{sheet_name}': Image at {image.anchor} has no alt text",
                            severity=Severity.SERIOUS,
                        ))

            wb.close()
            return issues
        except ImportError:
            return []
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze XLSX images {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]
