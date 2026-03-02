"""DOCX and PPTX accessibility checks."""
from pathlib import Path

from canvas_a11y.models import AccessibilityIssue, Severity
from canvas_a11y.checks.base import AccessibilityCheck
from canvas_a11y.checks.registry import register_check


@register_check
class DocxImagesMissingAlt(AccessibilityCheck):
    check_id = "docx-images-missing-alt"
    title = "DOCX images without alt text"
    description = "Images in Word documents should have alt text"
    wcag_criterion = "1.1.1"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".docx",):
            return []
        try:
            from docx import Document
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            import xml.etree.ElementTree as ET

            doc = Document(str(file_path))
            issues = []

            # Check inline images for alt text
            nsmap = {
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            }

            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    drawing_elements = run._element.findall('.//wp:inline', nsmap)
                    drawing_elements += run._element.findall('.//wp:anchor', nsmap)
                    for drawing in drawing_elements:
                        docPr = drawing.find('.//wp:docPr', nsmap)
                        if docPr is None:
                            # Try without namespace
                            for child in drawing.iter():
                                if child.tag.endswith('}docPr') or child.tag == 'docPr':
                                    docPr = child
                                    break

                        alt_text = ""
                        if docPr is not None:
                            alt_text = docPr.get('descr', '') or docPr.get('title', '')

                        if not alt_text.strip():
                            name = docPr.get('name', 'unknown') if docPr is not None else 'unknown'
                            issues.append(self._make_issue(
                                f"Image '{name}' in DOCX has no alt text",
                                severity=Severity.SERIOUS,
                            ))

            return issues
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze DOCX {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]


@register_check
class PptxSlidesMissingTitles(AccessibilityCheck):
    check_id = "pptx-slides-missing-titles"
    title = "PPTX slides without titles"
    description = "PowerPoint slides should have titles for navigation"
    wcag_criterion = "2.4.2"

    def check_html(self, html: str, url: str = "") -> list[AccessibilityIssue]:
        return []

    def check_file(self, file_path: Path) -> list[AccessibilityIssue]:
        if file_path.suffix.lower() not in (".pptx",):
            return []
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            issues = []

            for i, slide in enumerate(prs.slides, 1):
                has_title = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.shape_id == slide.shapes.title_shape_id if hasattr(slide.shapes, 'title_shape_id') else False:
                        has_title = True
                        break
                    if hasattr(shape, "name") and "title" in shape.name.lower():
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_title = True
                            break

                # Simpler check: just look for slide.shapes.title
                if not has_title and slide.shapes.title is not None:
                    if slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text.strip():
                        has_title = True

                if not has_title:
                    issues.append(self._make_issue(
                        f"Slide {i} has no title",
                        severity=Severity.SERIOUS,
                    ))

            return issues
        except Exception as e:
            return [self._make_issue(
                f"Could not analyze PPTX {file_path.name}: {e}",
                severity=Severity.MODERATE,
            )]
